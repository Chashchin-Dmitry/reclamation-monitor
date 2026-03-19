"""
Улучшенный классификатор рекламаций на Python

Этот модуль содержит классы для:
1. Классификации рекламаций на основе темы, содержимого письма и вложений
2. Обработки вложений разных типов
"""
import ast
import os
import pandas as pd
import re
import logging
import json
import email as email_module
from email.header import decode_header as decode_header_util
from pathlib import Path
from typing import Dict, List, Optional, Any, Tuple, Union

# Настройка логирования
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler("reclamation_classifier.log"),
        logging.StreamHandler()
    ]
)
logger = logging.getLogger("ReclamationClassifier")


class ReclamationClassifier:
    """
    Классификатор рекламаций v2.0

    ВАЖНО: Этот классификатор НЕ определяет is_reclamation!
    is_reclamation определяется ТОЛЬКО LLaMA.

    Классификатор:
    - Проверяет blacklist (счета, бухгалтерия)
    - Определяет категорию (Наземка/Метро/ЖДТ/Спецтехника)
    - Вычисляет score для уверенности в категории
    """

    # Путь к файлу конфигурации
    CONFIG_FILE = "classifier_config.json"

    def __init__(self):
        # Загружаем конфигурацию из файла
        self._config = self._load_config()

        # Основные типы рекламаций
        self.reclamation_types = [
            'Рекламации Наземка',
            'Рекламации Метро',
            'Рекламации Спецтехника',
            'Рекламации ЖДТ'
        ]

        # Подкатегории из конфига
        self.subcategories = self._config.get('subcategories', {
            'Метро': ['ЦТОВ', 'ТМХ', 'МВМ', 'ГУП Метрополитен'],
            'Спецтехника': ['Автоком', 'Белаз', 'Пожтехпром'],
        })

        # Ключевые слова для категорий (из конфига, НЕ хардкод)
        self.keywords_for_categories = self._config.get('keywords_for_categories', {})

        # Компании и их категории (из конфига, НЕ хардкод)
        self.company_to_category = self._config.get('company_to_category', {})

        # Blacklist ключевые слова (из конфига)
        self.blacklist_keywords = self._config.get('blacklist_keywords', [])

        # Мапинг для определения получателей
        self.reclamation_distribution_map = {}

        logger.info(f"[CLASSIFIER] Загружено: {len(self.blacklist_keywords)} blacklist, "
                   f"{len(self.keywords_for_categories)} категорий, "
                   f"{len(self.company_to_category)} компаний")

    def _load_config(self) -> Dict[str, Any]:
        """Загружает конфигурацию из JSON файла."""
        try:
            config_path = Path(__file__).parent / self.CONFIG_FILE
            if config_path.exists():
                with open(config_path, 'r', encoding='utf-8') as f:
                    config = json.load(f)
                    logger.info(f"[CLASSIFIER] Конфиг загружен: {config_path}")
                    return config
            else:
                logger.warning(f"[CLASSIFIER] Конфиг не найден: {config_path}, используем defaults")
                return {}
        except Exception as e:
            logger.error(f"[CLASSIFIER] Ошибка загрузки конфига: {e}")
            return {}

    def is_blacklisted(self, email_data: Dict[str, Any]) -> Tuple[bool, Optional[str]]:
        """
        Проверяет, попадает ли письмо в blacklist (бухгалтерия, счета).

        Args:
            email_data: Данные письма

        Returns:
            (is_blacklisted, matched_keyword) — кортеж с результатом и найденным ключевым словом
        """
        try:
            subject = (email_data.get('subject') or '').lower()
            # Проверяем только тему и первые 500 символов тела (не всё письмо)
            body = (email_data.get('body') or '')[:500].lower()
            text = f"{subject} {body}"

            for keyword in self.blacklist_keywords:
                if keyword.lower() in text:
                    logger.info(f"[BLACKLIST] Совпадение: '{keyword}' в теме/теле письма")
                    return True, keyword

            return False, None

        except Exception as e:
            logger.warning(f"[BLACKLIST] Ошибка проверки: {e}")
            return False, None  # При ошибке — пропускаем дальше
    
    def load_distribution_map_from_file(self, distribution_file: str) -> None:
        """
        Загружает карту распределения рекламаций из файла
        
        Args:
            distribution_file: Путь к файлу с распределением
        """
        try:
            # Проверяем существование файла
            if not os.path.exists(distribution_file):
                logger.error(f"Файл распределения не найден: {distribution_file}")
                return
            
            # Определяем формат файла по расширению
            extension = os.path.splitext(distribution_file)[1].lower()
            
            if extension == '.csv':
                self._load_from_csv(distribution_file)
            elif extension in ['.xlsx', '.xls']:
                self._load_from_excel(distribution_file)
            else:
                logger.error(f"Неподдерживаемый формат файла распределения: {extension}")
        
        except Exception as e:
            logger.error(f"Ошибка при загрузке карты распределения: {e}")
            import traceback
            logger.error(traceback.format_exc())
        
    def _load_from_csv(self, csv_file: str) -> None:
        """
        Загружает карту распределения из CSV файла
        
        Args:
            csv_file: Путь к CSV файлу
        """
        import pandas as pd
        
        try:
            # Загружаем CSV
            df = pd.read_csv(csv_file, encoding='utf-8')
            self._process_distribution_dataframe(df)
        except Exception as e:
            logger.error(f"Ошибка при загрузке CSV файла: {e}")

    def _load_from_excel(self, excel_file: str) -> None:
        """
        Загружает карту распределения из Excel файла
        
        Args:
            excel_file: Путь к Excel файлу
        """
        import pandas as pd
        
        try:
            # Имя листа с данными о сотрудниках и их электронной почте
            sheet_name = "Сотр - имейл"
            
            logger.info(f"Загрузка данных из листа '{sheet_name}' файла {excel_file}")
            
            # Загружаем данные из конкретного листа Excel
            try:
                df = pd.read_excel(excel_file, sheet_name=sheet_name)
            except ValueError as e:
                if "No sheet named" in str(e):
                    # Если лист не найден, пробуем загрузить первый лист
                    logger.warning(f"Лист '{sheet_name}' не найден, используем первый лист")
                    df = pd.read_excel(excel_file)
                else:
                    raise
            
            # Проверяем наличие обязательных колонок
            required_columns = ['employee', 'email']
            missing_columns = [col for col in required_columns if col not in df.columns]
            
            if missing_columns:
                logger.error(f"В листе отсутствуют обязательные колонки: {missing_columns}")
                return
            
            # Проверяем наличие колонок с распределением
            distribution_columns = ['Рекламация_отправить_список', 'Рекламация_копия_список']
            missing_distr_columns = [col for col in distribution_columns if col not in df.columns]
            
            if missing_distr_columns:
                logger.warning(f"Отсутствуют колонки распределения: {missing_distr_columns}")
                # Пытаемся найти альтернативные колонки
                alternative_columns = []
                for col in df.columns:
                    if 'рекламац' in col.lower() or 'список' in col.lower():
                        alternative_columns.append(col)
                
                if alternative_columns:
                    logger.info(f"Найдены альтернативные колонки: {alternative_columns}")
                    
                    # Если найдена только одна колонка, используем её для обоих списков
                    if len(alternative_columns) == 1:
                        df['Рекламация_отправить_список'] = df[alternative_columns[0]]
                        df['Рекламация_копия_список'] = [[] for _ in range(len(df))]
                    # Если найдены две колонки, предполагаем, что первая - для получателей, вторая - для копий
                    elif len(alternative_columns) >= 2:
                        df['Рекламация_отправить_список'] = df[alternative_columns[0]]
                        df['Рекламация_копия_список'] = df[alternative_columns[1]]
            
            # Обрабатываем данные
            self._process_distribution_dataframe(df)
            
        except Exception as e:
            logger.error(f"Ошибка при загрузке Excel файла: {e}")
            import traceback
            logger.error(traceback.format_exc())

    def _process_distribution_dataframe(self, df):
        """
        Обрабатывает DataFrame с данными о распределении
        
        Args:
            df: Pandas DataFrame с данными
        """
        distribution_map = {}
        
        # Проверяем наличие необходимых колонок
        required_columns = ['employee', 'email', 'Рекламация_отправить_список', 'Рекламация_копия_список']
        missing_columns = [col for col in required_columns if col not in df.columns]
        
        if missing_columns:
            logger.error(f"В файле распределения отсутствуют колонки: {missing_columns}")
            return
        
        # Проходим по строкам DataFrame
        for _, row in df.iterrows():
            employee = row.get('employee', '')
            email_addr = row.get('email', '')
            send_list = row.get('Рекламация_отправить_список', '')
            copy_list = row.get('Рекламация_копия_список', '')
            
            if not email_addr or pd.isna(email_addr):
                continue
            
            # Преобразуем строки списков в фактические списки
            try:
                if isinstance(send_list, str) and send_list.startswith('[') and send_list.endswith(']'):
                    send_categories = ast.literal_eval(send_list)
                else:
                    send_categories = []

                if isinstance(copy_list, str) and copy_list.startswith('[') and copy_list.endswith(']'):
                    copy_categories = ast.literal_eval(copy_list)
                else:
                    copy_categories = []
            except Exception as e:
                logger.error(f"Ошибка при обработке списков для {employee}: {e}")
                send_categories = []
                copy_categories = []
            
            # Обрабатываем категории, на которые сотрудник подписан
            for category in send_categories:
                if isinstance(category, str):
                    category = category.strip().strip("'")
                    if category not in distribution_map:
                        distribution_map[category] = {'получатели': [], 'копия': []}
                    if email_addr not in distribution_map[category]['получатели']:
                        distribution_map[category]['получатели'].append(email_addr)
            
            # Обрабатываем категории, которые сотрудник получает в копии
            for category in copy_categories:
                if isinstance(category, str):
                    category = category.strip().strip("'")
                    if category not in distribution_map:
                        distribution_map[category] = {'получатели': [], 'копия': []}
                    if email_addr not in distribution_map[category]['копия']:
                        distribution_map[category]['копия'].append(email_addr)
        
        # Проверяем полноту карты распределения
        for reclamation_type in self.reclamation_types:
            if reclamation_type not in distribution_map:
                logger.warning(f"Категория '{reclamation_type}' отсутствует в карте распределения")
        
        # Выводим информацию о распределении
        logger.info(f"Карта распределения рекламаций загружена: {len(distribution_map)} категорий")
        for category, recipients in distribution_map.items():
            logger.info(f"Категория: {category}")
            logger.info(f"  Получатели: {recipients['получатели']}")
            logger.info(f"  Копия: {recipients['копия']}")
        
        self.reclamation_distribution_map = distribution_map
    
    def load_distribution_map_from_csv(self, csv_data: List[Dict[str, str]]) -> None:
        """
        Загружает карту распределения рекламаций из данных CSV
        
        Args:
            csv_data: Данные из CSV-файла
        """
        try:
            distribution_map = {}
            
            for row in csv_data:
                employee = row.get('employee', '')
                email_addr = row.get('email', '')
                send_list = row.get('Рекламация_отправить_список', '')
                copy_list = row.get('Рекламация_копия_список', '')
                
                if not email_addr:
                    continue
                
                # Преобразуем строки списков в фактические списки
                try:
                    if send_list and send_list.startswith('[') and send_list.endswith(']'):
                        send_categories = ast.literal_eval(send_list)
                    else:
                        send_categories = []

                    if copy_list and copy_list.startswith('[') and copy_list.endswith(']'):
                        copy_categories = ast.literal_eval(copy_list)
                    else:
                        copy_categories = []
                except Exception as e:
                    logger.error(f"Ошибка при обработке списков для {employee}: {e}")
                    send_categories = []
                    copy_categories = []
                
                # Обрабатываем категории, на которые сотрудник подписан
                for category in send_categories:
                    category = category.strip().strip("'")
                    if category not in distribution_map:
                        distribution_map[category] = {'получатели': [], 'копия': []}
                    if email_addr not in distribution_map[category]['получатели']:
                        distribution_map[category]['получатели'].append(email_addr)
                
                # Обрабатываем категории, которые сотрудник получает в копии
                for category in copy_categories:
                    category = category.strip().strip("'")
                    if category not in distribution_map:
                        distribution_map[category] = {'получатели': [], 'копия': []}
                    if email_addr not in distribution_map[category]['копия']:
                        distribution_map[category]['копия'].append(email_addr)
            
            self.reclamation_distribution_map = distribution_map
            logger.info(f"Успешно загружена карта распределения рекламаций: {len(distribution_map)} категорий")
            
        except Exception as e:
            logger.error(f"Ошибка при загрузке карты распределения: {e}")
    
    def set_test_distribution_map(self) -> None:
        """Устанавливает тестовую карту распределения для отладки"""
        self.reclamation_distribution_map = {
            'Рекламации Категория1': {
                'получатели': ['engineer1@your-company.ru', 'engineer2@your-company.ru'],
                'копия': []
            },
            'Рекламации Категория2': {
                'получатели': [],
                'копия': ['manager@your-company.ru', 'director@your-company.ru']
            },
            'Рекламации Категория3': {
                'получатели': ['engineer1@your-company.ru', 'engineer3@your-company.ru'],
                'копия': []
            },
            'Рекламации Категория4': {
                'получатели': ['engineer4@your-company.ru'],
                'копия': []
            }
        }
        logger.info("Установлена тестовая карта распределения рекламаций")
    
    def classify_reclamation(self, email_data: Dict[str, Any], attachments: List[Dict[str, Any]]) -> Dict[str, Any]:
        """
        Классифицирует КАТЕГОРИЮ рекламации (НЕ определяет is_reclamation!).

        ВАЖНО: is_reclamation ВСЕГДА None — решение принимает LLaMA.

        Args:
            email_data: Данные письма (тема, содержимое и т.д.)
            attachments: Информация о вложениях

        Returns:
            Dict с результатами классификации:
            - is_reclamation: None (ВСЕГДА — решает LLaMA)
            - blacklisted: bool (True если счёт/оплата)
            - blacklist_keyword: str | None (какое слово сработало)
            - category: str | None (Рекламации Наземка/Метро/ЖДТ/Спецтехника)
            - score: int (уверенность в категории)
            - subcategories: list
        """
        result = {
            'is_reclamation': None,  # ВСЕГДА None — решает LLaMA!
            'blacklisted': False,
            'blacklist_keyword': None,
            'category': None,
            'score': 0,
            'subcategories': [],
            'recipients': [],
            'copy_to': [],
            'error': None
        }

        try:
            # 1. BLACKLIST CHECK (быстрый выход для бухгалтерии)
            is_blacklisted, blacklist_keyword = self.is_blacklisted(email_data)
            result['blacklisted'] = is_blacklisted
            result['blacklist_keyword'] = blacklist_keyword

            if is_blacklisted:
                logger.info(f"[CLASSIFIER] Blacklisted: '{blacklist_keyword}'")
                # НЕ возвращаем сразу — продолжаем для категоризации (на случай если LLaMA решит иначе)

            # 2. Собираем тексты для анализа категории
            subject = (email_data.get('subject') or '')
            body = (email_data.get('body') or '')
            attachment_names = ' '.join([att.get('filename', '') for att in (attachments or [])])

            text_for_analysis = f"{subject} {body} {attachment_names}".lower()

            # 3. Скоринг по категориям (БЕЗ весов — просто подсчёт совпадений)
            category_scores = {
                'Наземка': 0,
                'Метро': 0,
                'Спецтехника': 0,
                'ЖДТ': 0
            }

            # 4. Поиск по ключевым словам (вес = 1)
            for category, keywords in self.keywords_for_categories.items():
                for keyword in keywords:
                    try:
                        pattern = re.escape(keyword.lower())
                        matches = re.findall(pattern, text_for_analysis, re.IGNORECASE)
                        category_scores[category] += len(matches)
                    except re.error as re_err:
                        logger.warning(f"[CLASSIFIER] Regex error for '{keyword}': {re_err}")

            # 5. Поиск по компаниям (вес = 2 — компании важнее для категоризации)
            for company, category in self.company_to_category.items():
                try:
                    pattern = re.escape(company.lower())
                    matches = re.findall(pattern, text_for_analysis, re.IGNORECASE)
                    category_scores[category] += len(matches) * 2
                except re.error as re_err:
                    logger.warning(f"[CLASSIFIER] Regex error for '{company}': {re_err}")

            # 6. Анализируем текст из вложений (если есть)
            attachment_text = ""
            for att in (attachments or []):
                try:
                    if att.get('extracted_text'):
                        attachment_text += ' ' + att['extracted_text'].lower()
                except Exception:
                    continue

            if attachment_text:
                for category, keywords in self.keywords_for_categories.items():
                    for keyword in keywords:
                        try:
                            pattern = re.escape(keyword.lower())
                            matches = re.findall(pattern, attachment_text, re.IGNORECASE)
                            category_scores[category] += len(matches)
                        except re.error:
                            continue

                for company, category in self.company_to_category.items():
                    try:
                        pattern = re.escape(company.lower())
                        matches = re.findall(pattern, attachment_text, re.IGNORECASE)
                        category_scores[category] += len(matches)
                    except re.error:
                        continue

            # 7. Определяем категорию с максимальным скором
            max_score = max(category_scores.values())
            result['score'] = max_score

            if max_score > 0:
                top_category = max(category_scores.items(), key=lambda x: x[1])[0]
                reclamation_type = f"Рекламации {top_category}"
                result['category'] = reclamation_type

                # 8. Получатели и копии
                distribution = self.reclamation_distribution_map.get(reclamation_type, {})
                result['recipients'] = distribution.get('получатели', [])
                result['copy_to'] = distribution.get('копия', [])

                # 9. Подкатегории
                if top_category in self.subcategories:
                    for subcat in self.subcategories[top_category]:
                        try:
                            pattern = r'\b' + re.escape(subcat.lower()) + r'\b'
                            if re.search(pattern, text_for_analysis, re.IGNORECASE):
                                result['subcategories'].append(subcat)
                        except re.error:
                            continue

                logger.info(f"[CLASSIFIER] Категория: {reclamation_type}, score={max_score}")
            else:
                logger.info(f"[CLASSIFIER] Категория не определена (score=0)")

            return result

        except Exception as e:
            logger.error(f"[CLASSIFIER] Ошибка классификации: {e}")
            result['error'] = str(e)
            return result

        finally:
            logger.debug(f"[CLASSIFIER] Результат: blacklisted={result['blacklisted']}, "
                        f"category={result['category']}, score={result['score']}")


class AttachmentProcessor:
    """Класс для обработки вложений различных форматов"""
    
    def __init__(self):
        self.logger = logging.getLogger("AttachmentProcessor")
        
        # Поддерживаемые расширения файлов и соответствующие методы
        self.supported_extensions = {
            '.pdf': self.extract_text_from_pdf,
            '.docx': self.extract_text_from_docx,
            '.doc': self.extract_text_from_docx,
            '.txt': self.extract_text_from_txt,
            '.csv': self.extract_text_from_csv,
            '.xlsx': self.extract_text_from_excel,
            '.xls': self.extract_text_from_excel,
            '.pptx': self.extract_text_from_pptx,
            '.eml': self.extract_text_from_eml
        }
    
    def process_attachment(self, file_path: str) -> str:
        """
        Обрабатывает вложение и извлекает текст в зависимости от типа файла
        
        Args:
            file_path: Путь к файлу
            
        Returns:
            Извлеченный текст или сообщение об ошибке
        """
        try:
            if not os.path.exists(file_path):
                self.logger.error(f"Файл не существует: {file_path}")
                return f"[Ошибка: файл не существует: {file_path}]"
            
            file_extension = os.path.splitext(file_path)[1].lower()
            
            if file_extension in self.supported_extensions:
                extract_method = self.supported_extensions[file_extension]
                return extract_method(file_path)
            else:
                return f"[Формат файла {file_extension} не поддерживается для извлечения текста]"
        
        except Exception as e:
            self.logger.error(f"Ошибка при обработке вложения {file_path}: {e}")
            return f"[Ошибка при обработке вложения: {str(e)}]"
    
    def extract_text_from_pdf(self, file_path: str) -> str:
        """
        Извлекает текст из PDF файла
        
        Args:
            file_path: Путь к PDF файлу
            
        Returns:
            Извлеченный текст
        """
        try:
            # Проверяем наличие PyMuPDF
            try:
                import fitz
            except ImportError:
                self.logger.warning("Библиотека PyMuPDF (fitz) не установлена")
                return f"[PDF файл: {os.path.basename(file_path)} - требуется PyMuPDF для извлечения текста]"
            
            full_text = ""
            raw_text = ""  # Только текст без заголовков "Страница X:"
            doc = fitz.open(file_path)

            for page_num in range(len(doc)):
                try:
                    page = doc.load_page(page_num)
                    text = page.get_text()
                    raw_text += text
                    full_text += f"Страница {page_num + 1}:\n{text}\n\n"
                except Exception as e:
                    self.logger.error(f"Ошибка при извлечении текста со страницы {page_num + 1}: {e}")
                    full_text += f"Страница {page_num + 1}: [Ошибка чтения страницы]\n\n"

            doc.close()

            # Fallback на OCR если PyMuPDF вернул пустой текст (PDF-скан)
            # Проверяем raw_text (без заголовков), а не full_text (с "Страница X:")
            if not raw_text.strip():
                self.logger.info(f"PDF скан обнаружен (0 символов текста), запускаем OCR: {file_path}")
                try:
                    from ocr_processor import extract_text_from_scanned_pdf
                    full_text = extract_text_from_scanned_pdf(file_path)
                except Exception as e:
                    self.logger.error(f"Ошибка OCR для {file_path}: {e}")

            return full_text
        
        except Exception as e:
            self.logger.error(f"Ошибка при извлечении текста из PDF: {e}")
            return f"[Ошибка при извлечении текста из PDF: {str(e)}]"
    
    def extract_text_from_docx(self, file_path: str) -> str:
        """
        Извлекает текст из DOCX/DOC файла
        
        Args:
            file_path: Путь к DOCX/DOC файлу
            
        Returns:
            Извлеченный текст
        """
        try:
            # Проверяем существование файла
            if not os.path.exists(file_path):
                self.logger.error(f"Файл не существует: {file_path}")
                return f"[Ошибка: файл не существует]"
            
            # Проверяем размер файла
            file_size = os.path.getsize(file_path)
            if file_size == 0:
                self.logger.error(f"Файл пуст: {file_path}")
                return f"[Ошибка: файл пуст]"
            
            # Проверяем расширение файла
            file_extension = os.path.splitext(file_path)[1].lower()
            
            # Для DOC файлов (старый формат)
            if file_extension == '.doc':
                try:
                    # Пробуем использовать antiword (если установлен)
                    import subprocess
                    try:
                        text = subprocess.check_output(['antiword', file_path]).decode('utf-8', errors='ignore')
                        if text.strip():
                            return text
                    except (subprocess.SubprocessError, FileNotFoundError):
                        self.logger.warning("antiword не установлен или вызвал ошибку")
                        return f"[DOC файл: {os.path.basename(file_path)} - требуется antiword для извлечения текста]"
                except Exception as doc_error:
                    self.logger.error(f"Ошибка при извлечении текста из DOC: {doc_error}")
                    return f"[Не удалось обработать DOC файл: {str(doc_error)}]"
            
            # Для DOCX используем python-docx или docx2txt
            try:
                # Сначала пробуем docx2txt как более надежный вариант
                try:
                    import docx2txt
                    text = docx2txt.process(file_path)
                    if text.strip():
                        return text
                except ImportError:
                    self.logger.info("Библиотека docx2txt не установлена, пробуем python-docx")
                    
                # Если docx2txt не сработал, пробуем python-docx
                import docx
                doc = docx.Document(file_path)
                full_text = []
                
                # Извлекаем текст из параграфов
                for para in doc.paragraphs:
                    if para.text.strip():  # Не добавляем пустые параграфы
                        full_text.append(para.text)
                
                # Извлекаем текст из таблиц
                for table in doc.tables:
                    for row in table.rows:
                        row_text = []
                        for cell in row.cells:
                            if cell.text.strip():
                                row_text.append(cell.text.strip())
                        if row_text:  # Не добавляем пустые строки
                            full_text.append(' | '.join(row_text))
                
                # Объединяем текст
                result = '\n'.join(full_text)
                
                # Проверяем, что извлекли текст
                if not result.strip():
                    self.logger.warning(f"Из файла {file_path} не удалось извлечь текст")
                    return f"[Файл не содержит текста]"
                
                return result
            
            except Exception as docx_error:
                error_str = str(docx_error)
                self.logger.error(f"Ошибка при обработке DOCX/DOC: {error_str}")
                return f"[Не удалось обработать файл {os.path.basename(file_path)}: {error_str}]"
        
        except Exception as e:
            error_message = f"Ошибка при извлечении текста из документа Word: {str(e)}"
            self.logger.error(error_message)
            return f"[{error_message}]"
    
    def extract_text_from_txt(self, file_path: str) -> str:
        """
        Извлекает текст из TXT файла
        
        Args:
            file_path: Путь к TXT файлу
            
        Returns:
            Извлеченный текст
        """
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        except UnicodeDecodeError:
            # Если не удалось прочитать с UTF-8, пробуем другие кодировки
            try:
                with open(file_path, 'r', encoding='latin-1') as file:
                    return file.read()
            except Exception as e:
                self.logger.error(f"Ошибка при извлечении текста из TXT: {e}")
                return f"[Ошибка при извлечении текста из TXT: {str(e)}]"
        except Exception as e:
            self.logger.error(f"Ошибка при извлечении текста из TXT: {e}")
            return f"[Ошибка при извлечении текста из TXT: {str(e)}]"
    
    def extract_text_from_csv(self, file_path: str) -> str:
        """
        Извлекает текст из CSV файла
        
        Args:
            file_path: Путь к CSV файлу
            
        Returns:
            Извлеченный текст
        """
        try:
            import pandas as pd
            df = pd.read_csv(file_path)
            return df.to_string(index=False)
        except Exception as e:
            self.logger.error(f"Ошибка при извлечении текста из CSV: {e}")
            return f"[Ошибка при извлечении текста из CSV: {str(e)}]"
    
    def extract_text_from_excel(self, file_path: str) -> str:
        """
        Извлекает текст из Excel файла
        
        Args:
            file_path: Путь к Excel файлу
            
        Returns:
            Извлеченный текст
        """
        try:
            import pandas as pd
            # Определяем engine по расширению файла
            if file_path.lower().endswith('.xls'):
                df = pd.read_excel(file_path, engine='xlrd')
            else:
                df = pd.read_excel(file_path, engine='openpyxl')
            return df.to_string(index=False)
        except Exception as e:
            self.logger.error(f"Ошибка при извлечении текста из Excel: {e}")
            return f"[Ошибка при извлечении текста из Excel: {str(e)}]"
    
    def extract_text_from_pptx(self, file_path: str) -> str:
        """
        Извлекает текст из PPTX файла
        
        Args:
            file_path: Путь к PPTX файлу
            
        Returns:
            Извлеченный текст
        """
        try:
            # Проверяем наличие python-pptx
            try:
                import pptx
            except ImportError:
                self.logger.warning("Библиотека python-pptx не установлена")
                return f"[PPTX файл: {os.path.basename(file_path)} - требуется python-pptx для извлечения текста]"
            
            prs = pptx.Presentation(file_path)
            text_runs = []
            
            for i, slide in enumerate(prs.slides):
                text_runs.append(f"=== Слайд {i+1} ===")
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(shape.text)
            
            return '\n'.join(text_runs)
        
        except Exception as e:
            self.logger.error(f"Ошибка при извлечении текста из PPTX: {e}")
            return f"[Ошибка при извлечении текста из PPTX: {str(e)}]"

    def extract_text_from_eml(self, file_path: str) -> str:
        """
        Извлекает текст из .eml файла (вложенное письмо)

        Args:
            file_path: Путь к .eml файлу

        Returns:
            Извлеченный текст (заголовки + тело)
        """
        try:
            with open(file_path, 'rb') as f:
                msg = email_module.message_from_bytes(f.read())

            parts = []

            # Извлекаем заголовки
            for header_name in ('Subject', 'From', 'To', 'Date'):
                raw_value = msg.get(header_name, '')
                if raw_value:
                    # Декодируем заголовок (может быть в base64/quoted-printable)
                    decoded_parts = decode_header_util(raw_value)
                    decoded_value = ''
                    for part_bytes, charset in decoded_parts:
                        if isinstance(part_bytes, bytes):
                            decoded_value += part_bytes.decode(charset or 'utf-8', errors='replace')
                        else:
                            decoded_value += part_bytes
                    parts.append(f"{header_name}: {decoded_value}")

            parts.append('')  # пустая строка между заголовками и телом

            # Извлекаем тело письма
            if msg.is_multipart():
                for part in msg.walk():
                    content_type = part.get_content_type()
                    # Рекурсивная обработка вложенных message/rfc822
                    if content_type == 'message/rfc822':
                        payload = part.get_payload()
                        if payload and isinstance(payload, list):
                            for nested_msg in payload:
                                parts.append("--- Вложенное письмо ---")
                                parts.append(self._extract_text_from_email_message(nested_msg))
                    elif content_type == 'text/plain':
                        charset = part.get_content_charset() or 'utf-8'
                        try:
                            body = part.get_payload(decode=True).decode(charset, errors='replace')
                            parts.append(body)
                        except Exception:
                            pass
                    elif content_type == 'text/html' and not any('text/plain' == p.get_content_type() for p in msg.walk() if p != part):
                        # HTML fallback только если нет text/plain
                        charset = part.get_content_charset() or 'utf-8'
                        try:
                            html = part.get_payload(decode=True).decode(charset, errors='replace')
                            # Простая очистка HTML тегов
                            text = re.sub(r'<[^>]+>', ' ', html)
                            text = re.sub(r'\s+', ' ', text).strip()
                            parts.append(text)
                        except Exception:
                            pass
            else:
                charset = msg.get_content_charset() or 'utf-8'
                try:
                    body = msg.get_payload(decode=True).decode(charset, errors='replace')
                    parts.append(body)
                except Exception:
                    pass

            return '\n'.join(parts)

        except Exception as e:
            self.logger.error(f"Ошибка при извлечении текста из EML: {e}")
            return f"[Ошибка при извлечении текста из EML: {str(e)}]"

    def _extract_text_from_email_message(self, msg) -> str:
        """Рекурсивно извлекает текст из email.message.Message объекта"""
        parts = []
        for header_name in ('Subject', 'From', 'To', 'Date'):
            raw_value = msg.get(header_name, '')
            if raw_value:
                decoded_parts = decode_header_util(raw_value)
                decoded_value = ''
                for part_bytes, charset in decoded_parts:
                    if isinstance(part_bytes, bytes):
                        decoded_value += part_bytes.decode(charset or 'utf-8', errors='replace')
                    else:
                        decoded_value += part_bytes
                parts.append(f"{header_name}: {decoded_value}")

        if msg.is_multipart():
            for part in msg.walk():
                if part.get_content_type() == 'text/plain':
                    charset = part.get_content_charset() or 'utf-8'
                    try:
                        body = part.get_payload(decode=True).decode(charset, errors='replace')
                        parts.append(body)
                    except Exception:
                        pass
        else:
            charset = msg.get_content_charset() or 'utf-8'
            try:
                body = msg.get_payload(decode=True).decode(charset, errors='replace')
                parts.append(body)
            except Exception:
                pass

        return '\n'.join(parts)